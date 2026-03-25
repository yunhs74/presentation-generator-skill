#!/usr/bin/env python3
"""
Design Enhancer — 프레젠테이션의 디자인을 자동으로 개선합니다.

- 도형 배치 최적화 (그리드 정렬, 균등 분배)
- 데이터 특성에 따른 최적 차트 유형 선택
- SmartArt 스타일 도형 그룹 생성
- 텍스트 오버플로 방지
- 색상 일관성 적용
- 시각적 밸런스 개선
"""

import math
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE


# ─── 색상 팔레트 ─────────────────────────────────────────────────────

DEFAULT_PALETTES = {
    "professional": [
        RGBColor(0x2C, 0x3E, 0x50),  # Dark Blue
        RGBColor(0x34, 0x98, 0xDB),  # Blue
        RGBColor(0x1A, 0xBC, 0x9C),  # Teal
        RGBColor(0xE7, 0x4C, 0x3C),  # Red
        RGBColor(0xF3, 0x9C, 0x12),  # Orange
        RGBColor(0x9B, 0x59, 0xB6),  # Purple
    ],
    "modern": [
        RGBColor(0x66, 0x73, 0x8E),  # Slate
        RGBColor(0x56, 0xB4, 0xE9),  # Sky Blue
        RGBColor(0x00, 0x96, 0x88),  # Teal
        RGBColor(0xFF, 0x6F, 0x61),  # Coral
        RGBColor(0xFF, 0xC1, 0x07),  # Amber
        RGBColor(0xAB, 0x47, 0xBC),  # Purple
    ],
    "minimal": [
        RGBColor(0x33, 0x33, 0x33),  # Dark Gray
        RGBColor(0x55, 0x55, 0x55),  # Medium Gray
        RGBColor(0x00, 0x7A, 0xCC),  # Accent Blue
        RGBColor(0x88, 0x88, 0x88),  # Light Gray
        RGBColor(0xBB, 0xBB, 0xBB),  # Lighter Gray
        RGBColor(0xDD, 0xDD, 0xDD),  # Very Light Gray
    ],
}


def get_palette(scheme="auto", theme_colors=None):
    """색상 팔레트를 반환합니다."""
    if scheme in DEFAULT_PALETTES:
        return DEFAULT_PALETTES[scheme]
    if theme_colors:
        palette = []
        for key, val in theme_colors.items():
            if val and len(val) == 6:
                try:
                    palette.append(RGBColor(
                        int(val[0:2], 16),
                        int(val[2:4], 16),
                        int(val[4:6], 16),
                    ))
                except ValueError:
                    continue
        if palette:
            return palette
    return DEFAULT_PALETTES["professional"]


# ─── 차트 유형 추천 ─────────────────────────────────────────────────

def recommend_chart_type(chart_data):
    """데이터 특성에 따라 최적의 차트 유형을 추천합니다."""
    categories = chart_data.get("categories", [])
    series = chart_data.get("series", [])
    requested_type = chart_data.get("type", "auto")

    if requested_type != "auto":
        type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "radar": XL_CHART_TYPE.RADAR_FILLED,
            "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
        }
        return type_map.get(requested_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    num_categories = len(categories)
    num_series = len(series)

    # 단일 시리즈 + 적은 카테고리 → 파이/도넛
    if num_series == 1 and num_categories <= 6:
        values = series[0].get("values", [])
        if all(v >= 0 for v in values):
            return XL_CHART_TYPE.DOUGHNUT if num_categories > 3 else XL_CHART_TYPE.PIE

    # 시계열 데이터 (년도, 분기, 월 등) → 라인 차트
    time_keywords = ["Q", "월", "년", "일", "week", "month", "year", "quarter"]
    is_time_series = any(
        any(kw in str(cat) for kw in time_keywords) for cat in categories
    )
    if is_time_series:
        return XL_CHART_TYPE.LINE_MARKERS

    # 다중 시리즈 비교 → 클러스터 컨럼
    if num_series >= 2:
        return XL_CHART_TYPE.COLUMN_CLUSTERED

    # 많은 카테고리 → 가로 바 차트
    if num_categories > 8:
        return XL_CHART_TYPE.BAR_CLUSTERED

    # 기본값
    return XL_CHART_TYPE.COLUMN_CLUSTERED


# ─── 도형 배치 최적화 ────────────────────────────────────────────────

def calculate_grid_positions(count, slide_width, slide_height, margin=0.5):
    """항목 수에 따라 최적의 그리드 위치를 계산합니다."""
    usable_width = slide_width - 2 * margin
    usable_height = slide_height - 2 * margin - 1.5  # 제목 영역 제외

    if count <= 3:
        cols = count
        rows = 1
    elif count <= 6:
        cols = 3
        rows = math.ceil(count / 3)
    elif count <= 9:
        cols = 3
        rows = math.ceil(count / 3)
    else:
        cols = 4
        rows = math.ceil(count / 4)

    item_width = (usable_width - (cols - 1) * 0.2) / cols
    item_height = (usable_height - (rows - 1) * 0.2) / rows

    positions = []
    for i in range(count):
        row = i // cols
        col = i % cols
        x = margin + col * (item_width + 0.2)
        y = margin + 1.5 + row * (item_height + 0.2)  # 1.5" 제목 영역
        positions.append({
            "left": x,
            "top": y,
            "width": item_width,
            "height": item_height,
        })

    return positions


def align_shapes_to_grid(shapes, slide_width=10, slide_height=7.5):
    """기존 도형들을 균일한 그리드에 정렬합니다."""
    positions = calculate_grid_positions(
        len(shapes), slide_width, slide_height
    )
    for shape, pos in zip(shapes, positions):
        shape.left = Inches(pos["left"])
        shape.top = Inches(pos["top"])
        shape.width = Inches(pos["width"])
        shape.height = Inches(pos["height"])


# ─── SmartArt 스타일 도형 생성 ──────────────────────────────────────────

def create_process_shapes(slide, items, palette, slide_width=10, slide_height=7.5):
    """프로세스 흐름(→) 도형을 생성합니다."""
    n = len(items)
    margin = 0.5
    arrow_width = 0.3
    usable_width = slide_width - 2 * margin
    top = slide_height * 0.4
    shape_height = Inches(1.2)

    total_arrows = max(0, n - 1)
    total_shape_width = usable_width - total_arrows * arrow_width
    shape_width = total_shape_width / n

    shapes_created = []
    for i, item in enumerate(items):
        x = margin + i * (shape_width + arrow_width)
        color = palette[i % len(palette)]

        # 둥근 사각형 추가
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(shape_width), shape_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = str(item)
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

        shapes_created.append(shape)

        # 화살표 추가 (마지막 항목 제외)
        if i < n - 1:
            arrow_x = x + shape_width + 0.02
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x), Inches(top + 0.3),
                Inches(arrow_width - 0.04), Inches(0.6),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
            arrow.line.fill.background()

    return shapes_created


def create_cycle_shapes(slide, items, palette, slide_width=10, slide_height=7.5):
    """순환 구조(⟲) 도형을 생성합니다."""
    n = len(items)
    center_x = slide_width / 2
    center_y = slide_height / 2 + 0.3
    radius = min(slide_width, slide_height) * 0.25
    shape_size = 1.4

    shapes_created = []
    for i, item in enumerate(items):
        angle = (2 * math.pi * i / n) - math.pi / 2
        x = center_x + radius * math.cos(angle) - shape_size / 2
        y = center_y + radius * math.sin(angle) - shape_size / 2
        color = palette[i % len(palette)]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(shape_size), Inches(shape_size),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(item)
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        shapes_created.append(shape)

    # 중심에 순환 화살표 아이콘(원형)
    center_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 0.4), Inches(center_y - 0.4),
        Inches(0.8), Inches(0.8),
    )
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
    center_shape.line.fill.background()
    tf = center_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "⟲"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    return shapes_created


def create_hierarchy_shapes(slide, items, palette, slide_width=10, slide_height=7.5):
    """계층 구조(▽) 도형을 생성합니다."""
    n = len(items)
    margin = 1.0
    usable_width = slide_width - 2 * margin
    top_start = 1.8
    row_height = 0.9
    gap = 0.15

    shapes_created = []
    for i, item in enumerate(items):
        # 각 레벨의 너비가 점점 넓어지는 피라미드 구조
        level_ratio = (i + 1) / n
        level_width = usable_width * level_ratio
        x = (slide_width - level_width) / 2
        y = top_start + i * (row_height + gap)
        color = palette[i % len(palette)]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(level_width), Inches(row_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(item)
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        shapes_created.append(shape)

    return shapes_created


def create_comparison_shapes(slide, items, palette, slide_width=10, slide_height=7.5):
    """비교 구조(좌우 대비) 도형을 생성합니다. 2개 항목에 최적."""
    margin = 0.6
    gap = 0.4
    usable_width = slide_width - 2 * margin - gap
    shape_width = usable_width / 2
    top = 2.0
    shape_height = slide_height - top - 0.8

    shapes_created = []
    for i, item in enumerate(items[:2]):
        x = margin + i * (shape_width + gap)
        color = palette[i % len(palette)]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(shape_width), Inches(shape_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = str(item)
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        shapes_created.append(shape)

    return shapes_created


SMARTART_CREATORS = {
    "process": create_process_shapes,
    "cycle": create_cycle_shapes,
    "hierarchy": create_hierarchy_shapes,
    "comparison": create_comparison_shapes,
}


def create_smartart(slide, smartart_type, items, palette, slide_width=10, slide_height=7.5):
    """SmartArt 스타일 도형 그룹을 생성합니다."""
    creator = SMARTART_CREATORS.get(smartart_type, create_process_shapes)
    return creator(slide, items, palette, slide_width, slide_height)


# ─── 텍스트 오버플로 방지 ───────────────────────────────────────────

def auto_fit_text(text_frame, max_font_size=24, min_font_size=8):
    """텍스트 길이에 따라 폰트 크기를 자동 조정합니다."""
    total_text = text_frame.text
    text_len = len(total_text)

    if text_len <= 20:
        target_size = max_font_size
    elif text_len <= 50:
        target_size = max(min_font_size, max_font_size - 4)
    elif text_len <= 100:
        target_size = max(min_font_size, max_font_size - 8)
    elif text_len <= 200:
        target_size = max(min_font_size, max_font_size - 12)
    else:
        target_size = min_font_size

    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(target_size)

    return target_size


# ─── 테이블 스타일링 ────────────────────────────────────────────────

def style_table(table, palette, header_color=None):
    """테이블에 스타일을 적용합니다."""
    if header_color is None:
        header_color = palette[0] if palette else RGBColor(0x2C, 0x3E, 0x50)

    # 헤더 행 스타일
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True
                run.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER

    # 데이터 행 교대 색상
    for row_idx in range(1, len(table.rows)):
        for cell in table.rows[row_idx].cells:
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# ─── 슬라이드 밸런스 최적화 ─────────────────────────────────────────

def optimize_slide_balance(slide, slide_width=10, slide_height=7.5):
    """슬라이드의 시각적 밸런스를 개선합니다."""
    shapes = [s for s in slide.shapes if not s.is_placeholder]
    if len(shapes) < 2:
        return

    # 도형들의 수직 중심 정렬
    shape_centers = []
    for shape in shapes:
        cx = (shape.left + shape.width / 2) / 914400
        cy = (shape.top + shape.height / 2) / 914400
        shape_centers.append((cx, cy))

    # 전체 중심에서의 편차 확인
    avg_cx = sum(c[0] for c in shape_centers) / len(shape_centers)
    slide_center = slide_width / 2

    offset = slide_center - avg_cx
    if abs(offset) > 0.3:  # 0.3인치 이상 편차인 경우 보정
        for shape in shapes:
            new_left = shape.left + Inches(offset)
            if Inches(0.3) <= new_left <= Inches(slide_width - 0.3):
                shape.left = int(new_left)
