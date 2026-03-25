#!/usr/bin/env python3
"""
Template Analyzer — PowerPoint 템플릿(.pptx)의 구조를 분석합니다.

슬라이드 레이아웃, 플레이스홀더, 도형, 색상, 폰트 등을 추출하여
JSON 형태로 출력합니다.
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx 패키지가 필요합니다.")
    print("설치: pip install python-pptx")
    sys.exit(1)


def emu_to_inches(emu_value):
    """EMU 값을 인치로 변환합니다."""
    if emu_value is None:
        return None
    return round(emu_value / 914400, 2)


def extract_color(color_obj):
    """색상 객체에서 RGB 값을 추출합니다."""
    try:
        if color_obj and color_obj.rgb:
            return str(color_obj.rgb)
    except (AttributeError, TypeError):
        pass
    try:
        if color_obj and color_obj.theme_color:
            return f"theme:{color_obj.theme_color}"
    except (AttributeError, TypeError):
        pass
    return None


def analyze_shape(shape):
    """개별 도형을 분석합니다."""
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else "unknown",
        "position": {
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
        },
        "size": {
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
        },
    }

    # 텍스트 프레임 분석
    if shape.has_text_frame:
        tf = shape.text_frame
        paragraphs = []
        for para in tf.paragraphs:
            para_info = {
                "text": para.text,
                "alignment": str(para.alignment) if para.alignment else None,
            }
            if para.runs:
                run = para.runs[0]
                font_info = {}
                if run.font.size:
                    font_info["size_pt"] = run.font.size.pt
                if run.font.name:
                    font_info["name"] = run.font.name
                if run.font.bold is not None:
                    font_info["bold"] = run.font.bold
                if run.font.italic is not None:
                    font_info["italic"] = run.font.italic
                font_info["color"] = extract_color(run.font.color)
                para_info["font"] = font_info
            paragraphs.append(para_info)
        info["text_content"] = paragraphs

    # 테이블 분석
    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cell_texts": [
                [cell.text for cell in row.cells] for row in table.rows
            ],
        }

    # 차트 분석
    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
        }

    # 그룹 도형 분석
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["group_shapes"] = [
            analyze_shape(child) for child in shape.shapes
        ]

    # 플레이스홀더 정보
    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["placeholder"] = {
            "idx": ph.idx,
            "type": str(ph.type) if ph.type else None,
        }

    return info


def analyze_slide_layout(layout):
    """슬라이드 레이아웃을 분석합니다."""
    placeholders = []
    for ph in layout.placeholders:
        placeholders.append({
            "idx": ph.placeholder_format.idx,
            "name": ph.name,
            "type": str(ph.placeholder_format.type) if ph.placeholder_format.type else None,
            "position": {
                "left": emu_to_inches(ph.left),
                "top": emu_to_inches(ph.top),
            },
            "size": {
                "width": emu_to_inches(ph.width),
                "height": emu_to_inches(ph.height),
            },
        })

    return {
        "name": layout.name,
        "placeholders": placeholders,
        "shape_count": len(layout.shapes),
    }


def analyze_slide(slide, index):
    """슬라이드를 분석합니다."""
    shapes = [analyze_shape(shape) for shape in slide.shapes]

    layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"

    return {
        "index": index,
        "layout": layout_name,
        "shapes": shapes,
        "shape_count": len(shapes),
    }


def analyze_theme_colors(prs):
    """테마 색상을 분석합니다."""
    colors = {}
    try:
        theme = prs.slide_masters[0].element
        # 테마 XML에서 색상 추출 시도
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        color_scheme = theme.findall(".//a:clrScheme/*", ns)
        for elem in color_scheme:
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            for child in elem:
                child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if child_tag == "srgbClr":
                    colors[tag] = child.get("val")
                elif child_tag == "sysClr":
                    colors[tag] = child.get("lastClr", child.get("val"))
    except Exception:
        pass
    return colors


def analyze_template(template_path):
    """전체 템플릿을 분석합니다."""
    prs = Presentation(template_path)

    # 슬라이드 크기
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    # 슬라이드 레이아웃 분석
    layouts = []
    for layout in prs.slide_layouts:
        layouts.append(analyze_slide_layout(layout))

    # 슬라이드 분석
    slides = []
    for i, slide in enumerate(prs.slides):
        slides.append(analyze_slide(slide, i))

    # 테마 색상
    theme_colors = analyze_theme_colors(prs)

    # 사용된 폰트 수집
    fonts = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)

    result = {
        "template_path": str(template_path),
        "slide_size": {
            "width_inches": slide_width,
            "height_inches": slide_height,
        },
        "slide_count": len(prs.slides),
        "available_layouts": layouts,
        "slides": slides,
        "theme_colors": theme_colors,
        "fonts_used": sorted(fonts),
    }

    return result


def main():
    parser = argparse.ArgumentParser(
        description="PowerPoint 템플릿(.pptx)의 구조를 분석합니다."
    )
    parser.add_argument(
        "template",
        help="분석할 템플릿 파일 경로 (.pptx)",
    )
    parser.add_argument(
        "--output", "-o",
        help="분석 결과를 저장할 JSON 파일 경로 (기본: 콘솔 출력)",
        default=None,
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        default=True,
        help="JSON 출력을 보기 좋게 포맷팅 (기본: 활성)",
    )

    args = parser.parse_args()

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"Error: 파일을 찾을 수 없습니다: {template_path}")
        sys.exit(1)
    if not template_path.suffix.lower() == ".pptx":
        print(f"Error: .pptx 파일만 지원됩니다: {template_path}")
        sys.exit(1)

    print(f"템플릿 분석 중: {template_path}")
    result = analyze_template(template_path)

    indent = 2 if args.pretty else None
    json_output = json.dumps(result, indent=indent, ensure_ascii=False)

    if args.output:
        output_path = Path(args.output)
        output_path.write_text(json_output, encoding="utf-8")
        print(f"분석 결과가 저장되었습니다: {output_path}")
    else:
        print("\n=== 분석 결과 ===\n")
        print(json_output)

    # 요약 출력
    print(f"\n=== 요약 ===")
    print(f"슬라이드 크기: {result['slide_size']['width_inches']}\" x {result['slide_size']['height_inches']}\"")
    print(f"슬라이드 수: {result['slide_count']}")
    print(f"사용 가능 레이아웃: {len(result['available_layouts'])}개")
    for layout in result["available_layouts"]:
        ph_names = [p["name"] for p in layout["placeholders"]]
        print(f"  - {layout['name']}: 플레이스홀더 {len(layout['placeholders'])}개 {ph_names}")
    print(f"사용된 폰트: {', '.join(result['fonts_used']) if result['fonts_used'] else '없음'}")


if __name__ == "__main__":
    main()
