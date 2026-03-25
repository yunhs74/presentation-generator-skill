#!/usr/bin/env python3
"""
Generate Presentation — 템플릿과 데이터를 기반으로 PowerPoint 프레젠테이션을 생성합니다.

사용법:
  python generate_presentation.py --template <template.pptx> --data <data.json> --output <output.pptx> [--enhance]
"""

import argparse
import json
import csv
import sys
import copy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
except ImportError:
    print("Error: python-pptx 패키지가 필요합니다.")
    print("설치: pip install python-pptx")
    sys.exit(1)

# 같은 디렉토리의 모듈 임포트
SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR))

from template_analyzer import analyze_template
from design_enhancer import (
    get_palette,
    recommend_chart_type,
    create_smartart,
    auto_fit_text,
    style_table,
    optimize_slide_balance,
    calculate_grid_positions,
)


# ─── 데이터 로더 ──────────────────────────────────────────────

def load_json_data(path):
    """JSON 파일에서 데이터를 로드합니다."""
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def load_csv_data(path):
    """CSV 파일에서 데이터를 로드하여 슬라이드 데이터로 변환합니다."""
    with open(path, "r", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        rows = list(reader)

    if not rows:
        return {"slides": []}

    headers = list(rows[0].keys())

    # CSV → 테이블 슬라이드로 변환
    return {
        "metadata": {"title": Path(path).stem},
        "slides": [
            {
                "layout": "Table",
                "content": {
                    "title": Path(path).stem,
                    "table": {
                        "headers": headers,
                        "rows": [[row.get(h, "") for h in headers] for row in rows],
                    },
                },
            }
        ],
    }


def load_excel_data(path):
    """Excel 파일에서 데이터를 로드합니다."""
    try:
        import openpyxl
    except ImportError:
        print("Error: openpyxl 패키지가 필요합니다.")
        print("설치: pip install openpyxl")
        sys.exit(1)

    wb = openpyxl.load_workbook(path, data_only=True)
    slides = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            str_row = [str(cell) if cell is not None else "" for cell in row]
            if i == 0:
                headers = str_row
            else:
                rows.append(str_row)

        if headers:
            slides.append({
                "layout": "Table",
                "content": {
                    "title": sheet_name,
                    "table": {
                        "headers": headers,
                        "rows": rows,
                    },
                },
            })

    return {
        "metadata": {"title": Path(path).stem},
        "slides": slides,
    }


def load_data(path):
    """파일 확장자에 따라 적절한 로더를 사용합니다."""
    ext = Path(path).suffix.lower()
    if ext == ".json":
        return load_json_data(path)
    elif ext == ".csv":
        return load_csv_data(path)
    elif ext in (".xlsx", ".xls"):
        return load_excel_data(path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} (JSON, CSV, Excel만 지원)")


# ─── 레이아웃 매칭 ─────────────────────────────────────────────

def find_best_layout(prs, requested_layout):
    """요청된 레이아웃 이름에 가장 적합한 슬라이드 레이아웃을 찾습니다."""
    requested_lower = requested_layout.lower().strip()

    # 정확한 이름 매칭
    for layout in prs.slide_layouts:
        if layout.name.lower().strip() == requested_lower:
            return layout

    # 부분 매칭
    for layout in prs.slide_layouts:
        if requested_lower in layout.name.lower() or layout.name.lower() in requested_lower:
            return layout

    # 키워드 기반 매칭
    keyword_map = {
        "title": ["title", "제목"],
        "content": ["content", "text", "본문", "내용"],
        "blank": ["blank", "빈", "empty"],
        "chart": ["chart", "차트", "graph"],
        "table": ["table", "표"],
        "two": ["two", "2", "column", "비교"],
        "picture": ["picture", "image", "사진", "이미지"],
        "smartart": ["smart", "process", "diagram"],
    }

    for layout in prs.slide_layouts:
        layout_lower = layout.name.lower()
        for key, keywords in keyword_map.items():
            if any(kw in requested_lower for kw in keywords):
                if any(kw in layout_lower for kw in keywords):
                    return layout

    # 기본 레이아웃 — 플레이스홀더가 가장 많은 레이아웃
    best = max(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))
    return best


# ─── 슬라이드 생성 ─────────────────────────────────────────────

def set_placeholder_text(slide, idx, text, font_size=None, bold=None, color=None, alignment=None):
    """플레이스홀더에 텍스트를 설정합니다."""
    try:
        ph = slide.placeholders[idx]
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(text)
        if font_size:
            p.font.size = Pt(font_size)
        if bold is not None:
            p.font.bold = bold
        if color:
            p.font.color.rgb = color
        if alignment:
            p.alignment = alignment
        return True
    except (KeyError, IndexError):
        return False


def add_title_slide(prs, slide_data, palette, enhance=False):
    """타이틀 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Title Slide"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    subtitle = content.get("subtitle", "")

    set_placeholder_text(slide, 0, title, font_size=32, bold=True)
    set_placeholder_text(slide, 1, subtitle, font_size=18)

    if enhance:
        # 타이틀 슬라이드에 악센트 라인 추가
        from pptx.enum.shapes import MSO_SHAPE
        slide_width = prs.slide_width / 914400
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(slide_width * 0.3), Inches(3.2),
            Inches(slide_width * 0.4), Inches(0.06),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = palette[1] if len(palette) > 1 else palette[0]
        line.line.fill.background()

    return slide


def add_content_slide(prs, slide_data, palette, enhance=False):
    """텍스트 콘텐츠 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Content"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    body = content.get("body", "")
    bullets = content.get("bullets", [])

    set_placeholder_text(slide, 0, title, font_size=28, bold=True)

    # 본문 또는 불릿 처리
    body_ph = None
    for idx in [1, 2, 10, 13, 14]:  # 일반적인 본문 플레이스홀더 인덱스
        try:
            body_ph = slide.placeholders[idx]
            break
        except KeyError:
            continue

    if body_ph:
        tf = body_ph.text_frame
        tf.clear()
        if bullets:
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = str(bullet)
                p.font.size = Pt(16)
                p.level = 0
                p.space_after = Pt(6)
        elif body:
            tf.paragraphs[0].text = str(body)
            if enhance:
                auto_fit_text(tf, max_font_size=18, min_font_size=10)

    if enhance and bullets and len(bullets) >= 3:
        # 불릿이 많으면 카드 레이아웃으로 개선 시도
        _try_card_layout(slide, bullets, palette, prs)

    return slide


def _try_card_layout(slide, bullets, palette, prs):
    """불릿 리스트를 카드 레이아웃으로 변환 시도합니다."""
    from pptx.enum.shapes import MSO_SHAPE

    if len(bullets) > 6:
        return  # 카드가 너무 많으면 불릿 유지

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    positions = calculate_grid_positions(len(bullets), slide_width, slide_height)

    for i, (bullet, pos) in enumerate(zip(bullets, positions)):
        color = palette[i % len(palette)]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(pos["left"]), Inches(pos["top"]),
            Inches(pos["width"]), Inches(pos["height"]),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = str(bullet)
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


def add_chart_slide(prs, slide_data, palette, enhance=False):
    """차트 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Blank"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    set_placeholder_text(slide, 0, title, font_size=28, bold=True)

    chart_info = content.get("chart", {})
    if not chart_info:
        return slide

    categories = chart_info.get("categories", [])
    series_list = chart_info.get("series", [])

    if not categories or not series_list:
        return slide

    # 차트 유형 결정
    if enhance:
        chart_type = recommend_chart_type(chart_info)
    else:
        type_str = chart_info.get("type", "column")
        type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
            "auto": XL_CHART_TYPE.COLUMN_CLUSTERED,
        }
        chart_type = type_map.get(type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # 차트 데이터 생성
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    # 차트 위치/크기
    slide_width = prs.slide_width / 914400
    x = Inches(0.8)
    y = Inches(1.8)
    cx = Inches(slide_width - 1.6)
    cy = Inches(4.5)

    chart_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = chart_frame.chart

    # 차트 스타일 적용
    if enhance:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        try:
            plot = chart.plots[0]
            for i, series in enumerate(plot.series):
                color = palette[i % len(palette)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
        except Exception:
            pass

    return slide


def add_table_slide(prs, slide_data, palette, enhance=False):
    """테이블 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Blank"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    set_placeholder_text(slide, 0, title, font_size=28, bold=True)

    table_info = content.get("table", {})
    if not table_info:
        return slide

    headers = table_info.get("headers", [])
    rows_data = table_info.get("rows", [])

    if not headers:
        return slide

    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)

    slide_width = prs.slide_width / 914400
    table_width = slide_width - 1.6
    col_width = table_width / num_cols

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.8),
        Inches(table_width), Inches(min(num_rows * 0.5, 5.0)),
    )
    table = table_shape.table

    # 헤더 입력
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)

    # 데이터 입력
    for row_idx, row_values in enumerate(rows_data):
        for col_idx, value in enumerate(row_values):
            if col_idx < num_cols:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)

    # 테이블 스타일 적용
    if enhance:
        style_table(table, palette)

    return slide


def add_smartart_slide(prs, slide_data, palette, enhance=False):
    """SmartArt 스타일 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Blank"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    set_placeholder_text(slide, 0, title, font_size=28, bold=True)

    smartart_info = content.get("smartart", {})
    if not smartart_info:
        return slide

    smartart_type = smartart_info.get("type", "process")
    items = smartart_info.get("items", [])

    if not items:
        return slide

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400

    create_smartart(slide, smartart_type, items, palette, slide_width, slide_height)

    return slide


def add_image_slide(prs, slide_data, palette, enhance=False):
    """이미지 슬라이드를 추가합니다."""
    content = slide_data.get("content", {})
    layout = find_best_layout(prs, slide_data.get("layout", "Blank"))
    slide = prs.slides.add_slide(layout)

    title = content.get("title", "")
    set_placeholder_text(slide, 0, title, font_size=28, bold=True)

    image_path = content.get("image_path", "")
    if image_path and Path(image_path).exists():
        slide_width = prs.slide_width / 914400
        slide.shapes.add_picture(
            image_path,
            Inches(1.0), Inches(1.8),
            Inches(slide_width - 2.0), Inches(5.0),
        )

    return slide


# ─── 메인 생성 로직 ─────────────────────────────────────────────

SLIDE_CREATORS = {
    "title": add_title_slide,
    "title slide": add_title_slide,
    "content": add_content_slide,
    "text": add_content_slide,
    "chart": add_chart_slide,
    "graph": add_chart_slide,
    "table": add_table_slide,
    "smartart": add_smartart_slide,
    "smart art": add_smartart_slide,
    "process": add_smartart_slide,
    "image": add_image_slide,
    "picture": add_image_slide,
}


def generate_presentation(template_path, data, output_path, enhance=False):
    """프레젠테이션을 생성합니다."""
    prs = Presentation(template_path)

    # 디자인 기본 설정
    design_prefs = data.get("design_preferences", {})
    color_scheme = design_prefs.get("color_scheme", "auto")

    # 템플릿에서 테마 색상 추출
    analysis = analyze_template(template_path)
    theme_colors = analysis.get("theme_colors", {})
    palette = get_palette(color_scheme, theme_colors)

    # enhance가 데이터에서도 설정 가능
    if design_prefs.get("enhance", False):
        enhance = True

    # 기존 슬라이드 보존 여부
    keep_existing = data.get("keep_existing_slides", False)
    if not keep_existing:
        # 기존 슬라이드 모두 제거
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get("r:id")
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

    # 슬라이드 생성
    slides_data = data.get("slides", [])
    for slide_data in slides_data:
        layout_type = slide_data.get("layout", "content").lower().strip()
        creator = SLIDE_CREATORS.get(layout_type, add_content_slide)

        # SmartArt 타입 특수 처리
        if layout_type == "smartart" or "smartart" in slide_data.get("content", {}):
            creator = add_smartart_slide

        slide = creator(prs, slide_data, palette, enhance)

        # 디자인 개선: 시각적 밸런스 최적화
        if enhance and slide:
            slide_width = prs.slide_width / 914400
            slide_height = prs.slide_height / 914400
            optimize_slide_balance(slide, slide_width, slide_height)

    # 메타데이터 설정
    metadata = data.get("metadata", {})
    if metadata.get("title"):
        prs.core_properties.title = metadata["title"]
    if metadata.get("author"):
        prs.core_properties.author = metadata["author"]

    # 저장
    prs.save(output_path)
    print(f"프레젠테이션이 생성되었습니다: {output_path}")
    print(f"  - 슬라이드 수: {len(prs.slides)}")
    print(f"  - 디자인 개선: {'활성' if enhance else '비활성'}")

    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="템플릿과 데이터를 기반으로 PowerPoint 프레젠테이션을 생성합니다."
    )
    parser.add_argument(
        "--template", "-t",
        required=True,
        help="템플릿 .pptx 파일 경로",
    )
    parser.add_argument(
        "--data", "-d",
        required=True,
        help="데이터 파일 경로 (JSON, CSV, Excel)",
    )
    parser.add_argument(
        "--output", "-o",
        default="output.pptx",
        help="출력 .pptx 파일 경로 (기본: output.pptx)",
    )
    parser.add_argument(
        "--enhance", "-e",
        action="store_true",
        help="디자인 자동 개선 활성화",
    )
    parser.add_argument(
        "--analyze-only",
        action="store_true",
        help="템플릿 분석만 수행 (프레젠테이션 미생성)",
    )

    args = parser.parse_args()

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"Error: 템플릿을 찾을 수 없습니다: {template_path}")
        sys.exit(1)

    if args.analyze_only:
        analysis = analyze_template(str(template_path))
        print(json.dumps(analysis, indent=2, ensure_ascii=False))
        return

    data_path = Path(args.data)
    if not data_path.exists():
        print(f"Error: 데이터 파일을 찾을 수 없습니다: {data_path}")
        sys.exit(1)

    data = load_data(str(data_path))
    generate_presentation(str(template_path), data, args.output, args.enhance)


if __name__ == "__main__":
    main()
